"""MCP Server for python-pptx — create and manipulate PowerPoint presentations."""

from __future__ import annotations

import asyncio
import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import Tool, TextContent

logger = logging.getLogger(__name__)

server = Server("pptx-server")

# ── In-memory state ──────────────────────────────────────────────
_prs: Presentation | None = None
_current_slide_idx: int = 0

# Standard layout names mapped to layout indices (0-based)
_STANDARD_LAYOUTS: dict[str, int] = {
    "title": 0,
    "title_and_content": 1,
    "section_header": 2,
    "two_content": 3,
    "comparison": 4,
    "title_only": 5,
    "blank": 6,
    "content_with_caption": 7,
    "picture_with_caption": 8,
}


def _ensure_prs() -> Presentation:
    if _prs is None:
        raise ValueError("No presentation created yet. Call create_presentation first.")
    return _prs


def _ensure_slide():
    prs = _ensure_prs()
    if len(prs.slides) == 0:
        raise ValueError("No slides in the presentation. Call add_slide first.")
    return prs.slides[_current_slide_idx], prs


# ── Tools ────────────────────────────────────────────────────────

@server.list_tools()
async def list_tools() -> list[Tool]:
    return [
        Tool(
            name="create_presentation",
            description="Create a new blank PowerPoint presentation in memory",
            inputSchema={"type": "object", "properties": {}, "required": []},
        ),
        Tool(
            name="add_slide",
            description="Add a new slide with a specified layout to the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "layout": {
                        "type": "string",
                        "enum": list(_STANDARD_LAYOUTS.keys()),
                        "description": "Slide layout type",
                        "default": "blank",
                    },
                },
                "required": [],
            },
        ),
        Tool(
            name="go_to_slide",
            description="Set the current working slide by index (1-based: 1 = first slide)",
            inputSchema={
                "type": "object",
                "properties": {
                    "index": {
                        "type": "integer",
                        "description": "1-based slide index to switch to",
                    },
                },
                "required": ["index"],
            },
        ),
        Tool(
            name="set_title",
            description="Set the title text of the current slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "Title text"},
                },
                "required": ["text"],
            },
        ),
        Tool(
            name="add_content",
            description="Add bullet points to the content placeholder of the current slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "bullets": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of bullet point strings",
                    },
                },
                "required": ["bullets"],
            },
        ),
        Tool(
            name="add_textbox",
            description="Add a freeform textbox to the current slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "Text content"},
                    "left_inches": {"type": "number", "description": "X position from left (inches)", "default": 1},
                    "top_inches": {"type": "number", "description": "Y position from top (inches)", "default": 2},
                    "width_inches": {"type": "number", "description": "Width (inches)", "default": 8},
                    "height_inches": {"type": "number", "description": "Height (inches)", "default": 3},
                    "font_size": {"type": "integer", "description": "Font size in points (e.g. 18)"},
                },
                "required": ["text"],
            },
        ),
        Tool(
            name="add_table",
            description="Add a table to the current slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "rows": {
                        "type": "array",
                        "items": {"type": "array", "items": {"type": "string"}},
                        "description": "2D array of cell values",
                    },
                    "left_inches": {"type": "number", "description": "X position (inches)", "default": 1},
                    "top_inches": {"type": "number", "description": "Y position (inches)", "default": 2},
                    "width_inches": {"type": "number", "description": "Total table width (inches)", "default": 8},
                    "height_inches": {"type": "number", "description": "Total table height (inches)", "default": 4},
                },
                "required": ["rows"],
            },
        ),
        Tool(
            name="add_picture",
            description="Insert an image onto the current slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute path to the image file"},
                    "left_inches": {"type": "number", "description": "X position (inches)", "default": 1},
                    "top_inches": {"type": "number", "description": "Y position (inches)", "default": 1},
                    "width_inches": {"type": "number", "description": "Width in inches (optional)"},
                    "height_inches": {"type": "number", "description": "Height in inches (optional)"},
                },
                "required": ["path"],
            },
        ),
        Tool(
            name="add_speaker_notes",
            description="Add speaker notes to the current slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "Speaker notes text"},
                },
                "required": ["text"],
            },
        ),
        Tool(
            name="get_presentation_info",
            description="Get info about the current presentation: slide count, current slide index",
            inputSchema={"type": "object", "properties": {}, "required": []},
        ),
        Tool(
            name="save_presentation",
            description="Save the presentation to a .pptx file on disk",
            inputSchema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute path for the output .pptx file"},
                },
                "required": ["path"],
            },
        ),
    ]


def _get_layout_idx(layout_name: str, prs: Presentation) -> int:
    """Map a layout name to the best available layout index."""
    idx = _STANDARD_LAYOUTS.get(layout_name, 6)  # default: blank
    total = len(prs.slide_layouts)
    return min(idx, max(0, total - 1))


@server.call_tool()
async def call_tool(name: str, arguments: dict) -> list[TextContent]:
    global _prs, _current_slide_idx

    try:
        if name == "create_presentation":
            _prs = Presentation()
            _current_slide_idx = 0
            return [TextContent(type="text", text="Blank presentation created.")]

        elif name == "add_slide":
            prs = _ensure_prs()
            layout_name = arguments.get("layout", "blank")
            layout_idx = _get_layout_idx(layout_name, prs)
            slide_layout = prs.slide_layouts[layout_idx]
            prs.slides.add_slide(slide_layout)
            _current_slide_idx = len(prs.slides) - 1
            return [TextContent(
                type="text",
                text=f"Added slide #{_current_slide_idx + 1} with layout '{layout_name}'.",
            )]

        elif name == "go_to_slide":
            prs = _ensure_prs()
            idx = arguments["index"] - 1  # 1-based → 0-based
            if idx < 0 or idx >= len(prs.slides):
                return [TextContent(type="text", text=f"Slide index out of range (1-{len(prs.slides)}).")]
            _current_slide_idx = idx
            return [TextContent(type="text", text=f"Current slide is now #{idx + 1}.")]

        elif name == "set_title":
            slide, prs = _ensure_slide()
            if slide.shapes.title:
                slide.shapes.title.text = arguments["text"]
            else:
                # Fallback: add a title textbox at the top
                left = Inches(1)
                top = Inches(0.5)
                width = Inches(8)
                height = Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = arguments["text"]
            return [TextContent(type="text", text=f"Title set: {arguments['text'][:80]}")]

        elif name == "add_content":
            slide, prs = _ensure_slide()
            bullets = arguments["bullets"]
            # Try to find the content placeholder (idx 1 is usually the body)
            body_shape = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    body_shape = shape
                    break
            if body_shape is None:
                # Fallback: add a textbox
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
                tf = txBox.text_frame
            else:
                tf = body_shape.text_frame

            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
            return [TextContent(type="text", text=f"Added {len(bullets)} bullet points.")]

        elif name == "add_textbox":
            slide, prs = _ensure_slide()
            left = Inches(arguments.get("left_inches", 1))
            top = Inches(arguments.get("top_inches", 2))
            width = Inches(arguments.get("width_inches", 8))
            height = Inches(arguments.get("height_inches", 3))
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = arguments["text"]
            if arguments.get("font_size"):
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(arguments["font_size"])
            return [TextContent(type="text", text=f"Added textbox: {arguments['text'][:80]}")]

        elif name == "add_table":
            slide, prs = _ensure_slide()
            rows = arguments["rows"]
            if not rows or not rows[0]:
                raise ValueError("rows must be a non-empty 2D array")
            nrows = len(rows)
            ncols = len(rows[0])
            left = Inches(arguments.get("left_inches", 1))
            top = Inches(arguments.get("top_inches", 2))
            width = Inches(arguments.get("width_inches", 8))
            height = Inches(arguments.get("height_inches", 4))
            table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
            tbl = table_shape.table
            for i, row_data in enumerate(rows):
                for j, cell_text in enumerate(row_data):
                    if j < ncols:
                        tbl.cell(i, j).text = str(cell_text)
            return [TextContent(type="text", text=f"Added table with {nrows} rows × {ncols} cols.")]

        elif name == "add_picture":
            slide, prs = _ensure_slide()
            path = arguments["path"]
            if not Path(path).is_file():
                return [TextContent(type="text", text=f"Error: file not found — {path}")]
            left = Inches(arguments.get("left_inches", 1))
            top = Inches(arguments.get("top_inches", 1))
            kwargs = {}
            if arguments.get("width_inches"):
                kwargs["width"] = Inches(arguments["width_inches"])
            if arguments.get("height_inches"):
                kwargs["height"] = Inches(arguments["height_inches"])
            slide.shapes.add_picture(path, left, top, **kwargs)
            return [TextContent(type="text", text=f"Added picture: {path}")]

        elif name == "add_speaker_notes":
            slide, prs = _ensure_slide()
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = arguments["text"]
            return [TextContent(type="text", text=f"Speaker notes added to slide #{_current_slide_idx + 1}.")]

        elif name == "get_presentation_info":
            prs = _ensure_prs()
            slide_count = len(prs.slides)
            layout_count = len(prs.slide_layouts)
            current = _current_slide_idx + 1 if slide_count > 0 else 0
            info_lines = [
                f"Slides: {slide_count}",
                f"Current slide: {current}",
                f"Available layouts: {layout_count}",
                f"Slide width: {prs.slide_width}",
                f"Slide height: {prs.slide_height}",
            ]
            return [TextContent(type="text", text="\n".join(info_lines))]

        elif name == "save_presentation":
            prs = _ensure_prs()
            path = arguments["path"]
            out = Path(path)
            out.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(out))
            size_kb = out.stat().st_size / 1024
            return [TextContent(type="text", text=f"Presentation saved: {path} ({size_kb:.1f} KB)")]

        else:
            return [TextContent(type="text", text=f"Unknown tool: {name}")]

    except ValueError as e:
        return [TextContent(type="text", text=str(e))]
    except Exception as e:
        logger.exception("Tool error")
        return [TextContent(type="text", text=f"Error: {e}")]


# ── Entry point ──────────────────────────────────────────────────

def main():
    logger.info("Starting pptx-server...")
    asyncio.run(_run())


async def _run():
    async with stdio_server() as (read_stream, write_stream):
        await server.run(read_stream, write_stream, server.create_initialization_options())


if __name__ == "__main__":
    main()
